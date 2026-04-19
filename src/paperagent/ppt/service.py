from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from paperagent.config import Settings
from paperagent.schemas.models import DeckOutline, SlideOutline
from paperagent.storage.repositories import ChunkRepository, PaperRepository


class PPTService:
    def __init__(self, settings: Settings, paper_repository: PaperRepository, chunk_repository: ChunkRepository) -> None:
        self.settings = settings
        self.paper_repository = paper_repository
        self.chunk_repository = chunk_repository

    def generate(self, paper_id: str, template_path: Path | None = None) -> dict:
        paper = self.paper_repository.get_paper(paper_id)
        if not paper:
            raise ValueError(f"Paper {paper_id} not found.")
        chunks = self.chunk_repository.list_chunks(paper_id)
        outline = self._build_outline(paper_id, paper.title, chunks)

        deck_dir = self.settings.deck_dir / paper_id
        deck_dir.mkdir(parents=True, exist_ok=True)
        outline_path = deck_dir / "deck.json"
        ppt_path = deck_dir / "paper_briefing.pptx"
        outline_path.write_text(self._outline_to_json(outline), encoding="utf-8")
        self._render_ppt(outline, ppt_path, template_path=template_path)
        return {
            "paper_id": paper_id,
            "title": outline.title,
            "ppt_path": str(ppt_path),
            "outline_path": str(outline_path),
            "slide_count": len(outline.slides),
        }

    def _build_outline(self, paper_id: str, title: str, chunks) -> DeckOutline:
        grouped = {
            "Background & Problem": self._select_chunks(chunks, ["abstract", "introduction", "background", "motivation"]),
            "Method": self._select_chunks(chunks, ["method", "approach", "model", "framework"]),
            "Experiments": self._select_chunks(chunks, ["experiment", "evaluation", "results", "dataset"]),
            "Conclusion & Limitations": self._select_chunks(chunks, ["conclusion", "discussion", "limitation", "future"]),
        }
        slides = [
            SlideOutline(
                title=title,
                bullets=[
                    "论文讲解与汇报辅助稿",
                    f"Paper ID: {paper_id}",
                    "本页用于快速介绍论文主题与定位",
                ],
                notes="Title slide",
            )
        ]
        for slide_title, selected in grouped.items():
            bullets = self._make_bullets(selected)
            if not bullets:
                bullets = ["当前没有匹配到明显的章节内容，可回到 CLI 中继续追问。"]
            slides.append(
                SlideOutline(
                    title=slide_title,
                    bullets=bullets[:5],
                    notes="Auto-generated from indexed chunks.",
                )
            )
        slides.append(
            SlideOutline(
                title="How To Present This Paper",
                bullets=[
                    "先讲论文解决了什么问题，再讲方法思路。",
                    "实验部分重点解释对比对象、数据集和关键结果。",
                    "最后补充局限性，帮助听众形成完整判断。",
                ],
                notes="Presentation guidance",
            )
        )
        return DeckOutline(paper_id=paper_id, title=title, slides=slides)

    def _select_chunks(self, chunks, keywords: list[str]):
        selected = []
        for chunk in chunks:
            text = f"{chunk.section_title}\n{chunk.content}".lower()
            if any(keyword in text for keyword in keywords):
                selected.append(chunk)
        return selected[:3]

    def _make_bullets(self, chunks) -> list[str]:
        bullets: list[str] = []
        for chunk in chunks:
            sentences = [piece.strip() for piece in chunk.content.replace("\n", " ").split(".") if piece.strip()]
            if sentences:
                bullets.append(f"{sentences[0][:160]} (p.{chunk.page_number})")
        return bullets

    def _render_ppt(self, outline: DeckOutline, output_path: Path, template_path: Path | None) -> None:
        presentation = Presentation(str(template_path)) if template_path else Presentation()
        for index, slide_outline in enumerate(outline.slides):
            if index == 0:
                self._add_title_slide(presentation, slide_outline)
            else:
                self._add_content_slide(presentation, slide_outline)
        presentation.save(str(output_path))

    def _add_title_slide(self, presentation: Presentation, slide_outline: SlideOutline) -> None:
        layout = presentation.slide_layouts[0] if len(presentation.slide_layouts) > 0 else presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = slide_outline.title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "\n".join(slide_outline.bullets)

    def _add_content_slide(self, presentation: Presentation, slide_outline: SlideOutline) -> None:
        layout = presentation.slide_layouts[1] if len(presentation.slide_layouts) > 1 else presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(layout)
        title = slide.shapes.title
        title.text = slide_outline.title
        if len(slide.placeholders) > 1:
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            for index, bullet in enumerate(slide_outline.bullets):
                paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
                paragraph.text = bullet
        else:
            text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(4.5))
            frame = text_box.text_frame
            for index, bullet in enumerate(slide_outline.bullets):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = bullet

    def _outline_to_json(self, outline: DeckOutline) -> str:
        payload = {
            "paper_id": outline.paper_id,
            "title": outline.title,
            "slides": [
                {"title": slide.title, "bullets": slide.bullets, "notes": slide.notes}
                for slide in outline.slides
            ],
        }
        return json.dumps(payload, ensure_ascii=False, indent=2)
