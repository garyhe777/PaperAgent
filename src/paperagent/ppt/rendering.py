from __future__ import annotations

import json
import subprocess
from pathlib import Path

from paperagent.config import Settings
from paperagent.schemas.models import RenderResult, SlideContent


class PPTRenderService:
    def __init__(self, settings: Settings) -> None:
        self.settings = settings

    def render(
        self,
        paper_id: str,
        deck_title: str,
        slides: list[SlideContent],
    ) -> RenderResult:
        deck_dir = self.settings.deck_dir / paper_id
        deck_dir.mkdir(parents=True, exist_ok=True)
        work_dir = deck_dir / "skill_builder"
        work_dir.mkdir(parents=True, exist_ok=True)

        content_path = work_dir / "deck_content.json"
        render_config_path = work_dir / "render_config.json"
        output_path = deck_dir / "output.pptx"
        content_path.write_text(
            json.dumps(
                {
                    "paper_id": paper_id,
                    "title": deck_title,
                    "slides": [
                        {
                            "type": slide.slide_type,
                            "title": slide.title,
                            "bullets": slide.bullets,
                            "notes": slide.notes,
                            "citations": slide.citations,
                            "layout_hint": slide.layout_hint,
                            "visual_intent": slide.visual_intent,
                        }
                        for slide in slides
                    ],
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )
        render_config_path.write_text(
            json.dumps(
                {
                    "paper_id": paper_id,
                    "title": deck_title,
                    "content_path": str(content_path),
                    "output_path": str(output_path),
                },
                ensure_ascii=False,
                indent=2,
            ),
            encoding="utf-8",
        )

        runtime_info = self._detect_runtime()
        if runtime_info["available"]:
            builder_script = Path(__file__).resolve().with_name("skill_builder.mjs")
            try:
                self._execute_builder(
                    builder_script=builder_script,
                    render_config_path=render_config_path,
                    work_dir=work_dir,
                )
                if not output_path.exists():
                    raise RuntimeError("Skill renderer finished without creating output.pptx.")
                return RenderResult(
                    ppt_path=str(output_path),
                    slide_count=len(slides),
                    renderer="skill",
                )
            except RuntimeError:
                pass

        self._render_with_python_pptx(output_path=output_path, deck_title=deck_title, slides=slides)
        return RenderResult(
            ppt_path=str(output_path),
            slide_count=len(slides),
            renderer="python-pptx",
        )

    def _detect_runtime(self) -> dict[str, bool]:
        command = [
            "node",
            "-e",
            "import('@oai/artifact-tool').then(()=>process.stdout.write('ok')).catch(()=>process.exit(1))",
        ]
        try:
            completed = subprocess.run(command, capture_output=True, text=True, check=False, timeout=20)
        except (OSError, subprocess.SubprocessError):
            return {"available": False}
        return {"available": completed.returncode == 0}

    def _execute_builder(self, builder_script: Path, render_config_path: Path, work_dir: Path) -> None:
        command = [
            "node",
            str(builder_script),
            str(render_config_path),
        ]
        completed = subprocess.run(
            command,
            cwd=str(work_dir),
            capture_output=True,
            text=True,
            check=False,
            timeout=120,
        )
        if completed.returncode != 0:
            raise RuntimeError(
                "Skill renderer failed. "
                + (completed.stderr.strip() or completed.stdout.strip() or "Unknown JS builder error.")
            )

    def _render_with_python_pptx(self, output_path: Path, deck_title: str, slides: list[SlideContent]) -> None:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except Exception as exc:  # noqa: BLE001
            raise RuntimeError(
                "PPT rendering failed: neither JS skill runtime nor python-pptx fallback is available."
            ) from exc

        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)

        for index, slide_spec in enumerate(slides):
            layout = presentation.slide_layouts[0] if index == 0 else presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(layout)

            title_shape = slide.shapes.title
            if title_shape is not None:
                title_shape.text = slide_spec.title or deck_title

            body_placeholder = None
            if len(slide.placeholders) > 1:
                body_placeholder = slide.placeholders[1]

            if body_placeholder is not None:
                text_frame = body_placeholder.text_frame
                text_frame.clear()
                bullets = slide_spec.bullets or [slide_spec.notes or "Summary slide"]
                for bullet_index, bullet in enumerate(bullets):
                    paragraph = text_frame.paragraphs[0] if bullet_index == 0 else text_frame.add_paragraph()
                    paragraph.text = bullet
                    paragraph.level = 0
                    for run in paragraph.runs:
                        run.font.size = Pt(20)
            elif slide_spec.notes:
                text_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.8))
                text_box.text_frame.text = slide_spec.notes

            if slide_spec.notes:
                slide.notes_slide.notes_text_frame.text = slide_spec.notes

        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
